import io
from datetime import datetime, date
from typing import List, Optional, Dict
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd


class PowerPointGenerator:
    def __init__(self, data_manager):
        self.data_manager = data_manager
    
    def create_project_presentation(self, selected_projects: List[str], start_date: date, end_date: date) -> Optional[bytes]:
        """Create PowerPoint presentation with project data and charts"""
        try:
            prs = Presentation()
            
            # Set slide dimensions for widescreen
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, selected_projects)
            
            # Project overview slide
            self._create_overview_slide(prs, selected_projects, start_date, end_date)
            
            # Individual project slides
            for project_name in selected_projects:
                self._create_project_detail_slide(prs, project_name, start_date, end_date)
            
            # Performance comparison slide
            if len(selected_projects) > 1:
                self._create_performance_comparison_slide(prs, selected_projects)
            
            # Gantt chart slide
            self._create_gantt_chart_slide(prs, selected_projects)
            
            # Financial dashboard slide
            self._create_financial_dashboard_slide(prs, selected_projects, start_date, end_date)
            
            # Summary and recommendations slide
            self._create_summary_slide(prs, selected_projects)
            
            # Save to bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            
            return pptx_buffer.getvalue()
            
        except Exception as e:
            print(f"Error creating PowerPoint presentation: {e}")
            return None
    
    def _create_title_slide(self, prs, selected_projects):
        """Create title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "تقرير إدارة المشاريع"
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        project_list = ', '.join(selected_projects) if len(selected_projects) <= 3 else f"{len(selected_projects)} مشروع"
        subtitle.text = f"شركة عبد الله السعيد للاستشارات الهندسية\n{project_list}\n{datetime.now().strftime('%Y-%m-%d')}"
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_overview_slide(self, prs, selected_projects, start_date, end_date):
        """Create project overview slide with key metrics"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "نظرة عامة على المشاريع"
        title_frame.paragraphs[0].font.size = Inches(0.5)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Calculate summary metrics
        total_budget = 0
        total_actual_cost = 0
        avg_completion = 0
        project_count = len(selected_projects)
        
        y_position = Inches(2)
        metrics_per_row = 2
        metric_width = Inches(5)
        metric_height = Inches(1)
        
        for i, project_name in enumerate(selected_projects):
            project_data = self.data_manager.get_project_by_name(project_name)
            progress_data = self.data_manager.get_progress_data(project_name)
            
            if project_data:
                total_budget += project_data.get('total_budget', 0)
            
            if not progress_data.empty:
                latest_progress = progress_data.iloc[-1]
                total_actual_cost += latest_progress.get('actual_cost', 0)
                avg_completion += latest_progress.get('actual_completion', 0)
            
            # Add individual project metric box
            col = i % metrics_per_row
            row = i // metrics_per_row
            x_pos = Inches(1 + col * 6)
            y_pos = y_position + row * Inches(1.5)
            
            # Project box
            project_box = slide.shapes.add_textbox(x_pos, y_pos, metric_width, metric_height)
            project_frame = project_box.text_frame
            
            budget = project_data.get('total_budget', 0) if project_data else 0
            actual_cost = latest_progress.get('actual_cost', 0) if not progress_data.empty else 0
            completion = latest_progress.get('actual_completion', 0) if not progress_data.empty else 0
            
            project_frame.text = f"{project_name}\nالميزانية: {budget:,.0f}\nالتكلفة الفعلية: {actual_cost:,.0f}\nنسبة الإنجاز: {completion:.1f}%"
        
        # Summary metrics at the bottom
        summary_y = y_position + ((project_count // metrics_per_row + 1) * Inches(1.5)) + Inches(1)
        
        avg_completion = avg_completion / project_count if project_count > 0 else 0
        
        summary_box = slide.shapes.add_textbox(Inches(2), summary_y, Inches(9), Inches(1.5))
        summary_frame = summary_box.text_frame
        summary_frame.text = f"إجمالي الميزانية: {total_budget:,.0f}\nإجمالي التكلفة الفعلية: {total_actual_cost:,.0f}\nمتوسط نسبة الإنجاز: {avg_completion:.1f}%\nعدد المشاريع: {project_count}"
        summary_frame.paragraphs[0].font.bold = True
        summary_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_project_detail_slide(self, prs, project_name, start_date, end_date):
        """Create detailed slide for individual project"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = f"تفاصيل المشروع: {project_name}"
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Project information
        project_data = self.data_manager.get_project_by_name(project_name)
        progress_data = self.data_manager.get_progress_data(project_name)
        
        if project_data:
            # Project details box
            info_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(3))
            info_frame = info_box.text_frame
            
            start_date_str = project_data['start_date'].strftime('%Y-%m-%d') if project_data['start_date'] else 'غير محدد'
            end_date_str = project_data['end_date'].strftime('%Y-%m-%d') if project_data['end_date'] else 'غير محدد'
            
            info_text = f"تاريخ البداية: {start_date_str}\n"
            info_text += f"تاريخ النهاية: {end_date_str}\n"
            info_text += f"إجمالي الميزانية: {project_data.get('total_budget', 0):,.0f}\n"
            info_text += f"الشركة المنفذة: {project_data.get('executing_company', '')}\n"
            info_text += f"نوع المشروع: {project_data.get('project_type', '')}"
            
            info_frame.text = info_text
        
        # Progress chart (simple bar chart)
        if not progress_data.empty:
            self._add_progress_chart_to_slide(slide, progress_data, Inches(7), Inches(2))
    
    def _create_performance_comparison_slide(self, prs, selected_projects):
        """Create performance comparison slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "مقارنة أداء المشاريع"
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Create comparison chart
        self._add_comparison_chart_to_slide(slide, selected_projects)
    
    def _create_gantt_chart_slide(self, prs, selected_projects):
        """Create Gantt chart slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "مخطط جانت للمشاريع"
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Simple Gantt representation using text and bars
        y_start = Inches(2)
        bar_height = Inches(0.4)
        project_spacing = Inches(0.8)
        
        for i, project_name in enumerate(selected_projects):
            project_data = self.data_manager.get_project_by_name(project_name)
            if project_data:
                y_pos = y_start + i * project_spacing
                
                # Project name
                name_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(3), bar_height)
                name_frame = name_box.text_frame
                name_frame.text = project_name
                name_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
                # Timeline bar (simplified representation)
                start_date = project_data.get('start_date')
                end_date = project_data.get('end_date')
                
                if start_date and end_date:
                    # Calculate bar length based on project duration
                    duration = (end_date - start_date).days
                    max_duration = 365  # Normalize to one year
                    bar_width = min(Inches(6), Inches(6 * duration / max_duration))
                    
                    # Add colored rectangle for timeline
                    timeline_box = slide.shapes.add_textbox(Inches(5), y_pos, bar_width, bar_height)
                    timeline_frame = timeline_box.text_frame
                    timeline_frame.text = f"{start_date.strftime('%m/%d')} - {end_date.strftime('%m/%d')}"
    
    def _create_financial_dashboard_slide(self, prs, selected_projects, start_date, end_date):
        """Create financial dashboard slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "لوحة المراقبة المالية"
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Financial metrics
        total_planned_cost = 0
        total_actual_cost = 0
        total_budget = 0
        
        for project_name in selected_projects:
            project_data = self.data_manager.get_project_by_name(project_name)
            progress_data = self.data_manager.get_progress_data(project_name)
            
            if project_data:
                total_budget += project_data.get('total_budget', 0)
            
            if not progress_data.empty:
                latest_progress = progress_data.iloc[-1]
                total_planned_cost += latest_progress.get('planned_cost', 0)
                total_actual_cost += latest_progress.get('actual_cost', 0)
        
        # Financial summary
        financial_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
        financial_frame = financial_box.text_frame
        
        cost_variance = total_actual_cost - total_planned_cost
        budget_usage = (total_actual_cost / total_budget * 100) if total_budget > 0 else 0
        
        financial_text = f"إجمالي الميزانية: {total_budget:,.0f}\n"
        financial_text += f"التكلفة المخططة: {total_planned_cost:,.0f}\n"
        financial_text += f"التكلفة الفعلية: {total_actual_cost:,.0f}\n"
        financial_text += f"انحراف التكلفة: {cost_variance:,.0f}\n"
        financial_text += f"نسبة استخدام الميزانية: {budget_usage:.1f}%"
        
        financial_frame.text = financial_text
        financial_frame.paragraphs[0].font.size = Inches(0.25)
        financial_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_summary_slide(self, prs, selected_projects):
        """Create summary and recommendations slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "الملخص والتوصيات"
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Generate automatic recommendations based on data
        recommendations = self._generate_recommendations(selected_projects)
        
        summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(4))
        summary_frame = summary_box.text_frame
        summary_frame.text = recommendations
        summary_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def _add_progress_chart_to_slide(self, slide, progress_data, x, y):
        """Add progress chart to slide"""
        # Create simple chart data
        chart_data = CategoryChartData()
        chart_data.categories = ['المخطط', 'الفعلي']
        
        latest_progress = progress_data.iloc[-1]
        chart_data.add_series('النسبة المئوية', [
            latest_progress.get('planned_completion', 0),
            latest_progress.get('actual_completion', 0)
        ])
        
        # Add chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, Inches(5), Inches(3),
            chart_data
        )
    
    def _add_comparison_chart_to_slide(self, slide, selected_projects):
        """Add comparison chart to slide"""
        chart_data = CategoryChartData()
        chart_data.categories = selected_projects
        
        planned_values = []
        actual_values = []
        
        for project_name in selected_projects:
            progress_data = self.data_manager.get_progress_data(project_name)
            if not progress_data.empty:
                latest_progress = progress_data.iloc[-1]
                planned_values.append(latest_progress.get('planned_completion', 0))
                actual_values.append(latest_progress.get('actual_completion', 0))
            else:
                planned_values.append(0)
                actual_values.append(0)
        
        chart_data.add_series('المخطط', planned_values)
        chart_data.add_series('الفعلي', actual_values)
        
        # Add chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(2), Inches(2), Inches(9), Inches(4),
            chart_data
        )
    
    def _generate_recommendations(self, selected_projects):
        """Generate automatic recommendations based on project data"""
        recommendations = "التوصيات الآلية:\n\n"
        
        for project_name in selected_projects:
            progress_data = self.data_manager.get_progress_data(project_name)
            project_data = self.data_manager.get_project_by_name(project_name)
            
            if not progress_data.empty and project_data:
                latest_progress = progress_data.iloc[-1]
                planned = latest_progress.get('planned_completion', 0)
                actual = latest_progress.get('actual_completion', 0)
                
                if actual < planned:
                    recommendations += f"• {project_name}: المشروع متأخر عن الجدولة المخططة ({actual:.1f}% مقابل {planned:.1f}%)\n"
                elif actual > planned:
                    recommendations += f"• {project_name}: المشروع متقدم عن الجدولة المخططة ({actual:.1f}% مقابل {planned:.1f}%)\n"
                else:
                    recommendations += f"• {project_name}: المشروع يسير وفق الخطة المحددة\n"
        
        recommendations += "\nالتوصيات العامة:\n"
        recommendations += "• مراجعة دورية للمشاريع المتأخرة\n"
        recommendations += "• تحسين تخصيص الموارد\n"
        recommendations += "• متابعة التكاليف الفعلية مقارنة بالمخططة"
        
        return recommendations